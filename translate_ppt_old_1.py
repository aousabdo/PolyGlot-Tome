import os
import argparse
import logging
import asyncio
import aiohttp
import re
import time
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List
from collections import defaultdict
from asyncio import Semaphore

# Import your existing translation functions and configurations
# Assuming your translation script is named 'translation_script.py'
from book_translator import (
    load_config,
    RateLimiter,
    translate_chunk,
    google_gemini_translate,
    openai_translate,
    claude_sonnet_translate,
    rate_limiters,
    config,
    preprocess_arabic_text
)

# Configure logging
logging.basicConfig(level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s')

# Load configuration
config = load_config()

# Set up API keys
OPENAI_API_KEY = config.get('OPENAI_API_KEY')
GOOGLE_GEMINI_API_KEY = config.get('GOOGLE_GEMINI_API_KEY')
ANTHROPIC_API_KEY = config.get('ANTHROPIC_API_KEY')

# Validate API keys
if not OPENAI_API_KEY:
    logging.error("OpenAI API key is missing.")
if not GOOGLE_GEMINI_API_KEY:
    logging.warning("Google Gemini API key is missing.")
if not ANTHROPIC_API_KEY:
    logging.warning("Claude Sonnet API key is missing.")

# Translation function wrapper
async def translate_text(text: str, source_lang: str, target_lang: str, model: str, retries: int = 3) -> str:
    # Wrap your translate_chunk function to match the expected parameters
    async with aiohttp.ClientSession() as session:
        return await translate_chunk(session, text, source_lang, target_lang, model, retries)

# New function to process text frames
async def process_text_frame(text_frame, source_lang, target_lang, model, retries):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text.strip()
            if original_text:
                logging.info(f"Translating text: {original_text[:50]}...")
                translated_text = await translate_text(original_text, source_lang, target_lang, model, retries)
                run.text = translated_text

# Updated process_shapes function
async def process_shapes(shapes, source_lang: str, target_lang: str, model: str, retries: int):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            await process_shapes(shape.shapes, source_lang, target_lang, model, retries)
        elif shape.has_text_frame:
            await process_text_frame(shape.text_frame, source_lang, target_lang, model, retries)
        elif shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        await process_text_frame(cell.text_frame, source_lang, target_lang, model, retries)
        # Handle other shape types if necessary (e.g., charts, SmartArt)

# Main translation function
async def translate_pptx(input_file: str, output_file: str, source_lang: str = 'Arabic', target_lang: str = 'English',
                         model: str = 'google_gemini', retries: int = 3):
    prs = Presentation(input_file)
    tasks = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        logging.info(f"Processing Slide {slide_number}")
        # Process shapes in the slide
        await process_shapes(slide.shapes, source_lang, target_lang, model, retries)
        # Introduce a delay if necessary to comply with rate limits
        await asyncio.sleep(0.1)

    prs.save(output_file)
    logging.info(f"Translation completed. Output saved to {output_file}")

def main():
    parser = argparse.ArgumentParser(description='Translate PowerPoint presentations while preserving formatting.')
    parser.add_argument('pptx_path', help='Path to the input PPTX file.')
    parser.add_argument('--source_lang', default='Arabic', help='Source language (default: Arabic).')
    parser.add_argument('--target_lang', default='English', help='Target language (default: English).')
    parser.add_argument('--model', choices=['google_gemini', 'openai', 'claude_sonnet'], default='openai',
                        help='LLM model to use for translation.')
    parser.add_argument('--output_path', help='Path to save the translated PPTX file.')
    parser.add_argument('--retries', type=int, default=3, help='Number of retries for failed API calls.')
    args = parser.parse_args()

    if not args.output_path:
        base_name = os.path.splitext(os.path.basename(args.pptx_path))[0]
        args.output_path = f'{base_name}_translated.pptx'

    if not os.path.exists(args.pptx_path):
        logging.error(f"The file {args.pptx_path} does not exist.")
        return

    asyncio.run(translate_pptx(args.pptx_path, args.output_path, args.source_lang, args.target_lang,
                               args.model, args.retries))

if __name__ == '__main__':
    main()
